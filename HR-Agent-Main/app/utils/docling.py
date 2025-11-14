"""
Document processing using Docling for structure-preserving extraction.
Supports PDF, DOCX, XLSX, PPTX with table and layout preservation.
"""

from typing import Dict, Any, Optional
from pathlib import Path
import tempfile
from app.core.config import settings
from app.core.logging import get_logger

logger = get_logger(__name__)


class DoclingProcessor:
    """
    Document processor using Docling for structure-preserving extraction.

    Features:
    - Multi-format support (PDF, DOCX, XLSX, PPTX, TXT, MD)
    - Table preservation
    - Layout-aware extraction
    - OCR support (optional)
    """

    def __init__(
        self,
        max_file_size_mb: Optional[int] = None,
        ocr_enabled: Optional[bool] = None,
        preserve_tables: Optional[bool] = None,
    ):
        """
        Initialize Docling processor.

        Args:
            max_file_size_mb: Maximum file size in MB (default from settings)
            ocr_enabled: Enable OCR for scanned PDFs (default from settings)
            preserve_tables: Preserve table structures (default from settings)
        """
        self.max_file_size_mb = max_file_size_mb or settings.docling_max_file_size_mb
        self.ocr_enabled = ocr_enabled if ocr_enabled is not None else settings.docling_ocr_enabled
        self.preserve_tables = preserve_tables if preserve_tables is not None else settings.docling_preserve_tables
        self.supported_formats = settings.docling_supported_formats_list

        logger.info(
            f"Docling processor initialized: max_size={self.max_file_size_mb}MB, "
            f"ocr={self.ocr_enabled}, tables={self.preserve_tables}"
        )

    async def process_file(self, file_path: Path) -> Dict[str, Any]:
        """
        Process a document file and extract structured content.

        Args:
            file_path: Path to the document file

        Returns:
            Dictionary with extracted content and metadata:
            {
                "content": str,  # Extracted text content
                "metadata": {
                    "title": str,
                    "page_count": int,
                    "has_tables": bool,
                    "file_type": str,
                    "tables": List[Dict],  # Extracted tables
                    "structure": Dict,  # Document structure
                }
            }
        """
        try:
            file_ext = file_path.suffix.lower().lstrip(".")

            # Validate file type
            if file_ext not in self.supported_formats:
                raise ValueError(
                    f"Unsupported file type: {file_ext}. "
                    f"Supported formats: {', '.join(self.supported_formats)}"
                )

            # Validate file size
            file_size_mb = file_path.stat().st_size / (1024 * 1024)
            if file_size_mb > self.max_file_size_mb:
                raise ValueError(
                    f"File size ({file_size_mb:.2f}MB) exceeds maximum ({self.max_file_size_mb}MB)"
                )

            logger.info(f"Processing file: {file_path.name} ({file_ext}, {file_size_mb:.2f}MB)")

            # Route to appropriate processor
            if file_ext in ["pdf"]:
                result = await self._process_pdf(file_path)
            elif file_ext in ["docx", "doc"]:
                result = await self._process_docx(file_path)
            elif file_ext in ["xlsx", "xls"]:
                result = await self._process_excel(file_path)
            elif file_ext in ["pptx", "ppt"]:
                result = await self._process_ppt(file_path)
            elif file_ext in ["txt", "md"]:
                result = await self._process_text(file_path)
            else:
                raise ValueError(f"Processor not implemented for: {file_ext}")

            logger.info(
                f"Successfully processed {file_path.name}: "
                f"{len(result['content'])} chars, "
                f"{result['metadata'].get('page_count', 0)} pages"
            )

            return result

        except Exception as e:
            logger.error(f"Error processing file {file_path}: {e}", exc_info=True)
            raise

    async def _process_pdf(self, file_path: Path) -> Dict[str, Any]:
        """Process PDF using Docling."""
        try:
            from docling.document_converter import DocumentConverter

            converter = DocumentConverter()
            result = converter.convert(str(file_path))

            # Extract content and metadata
            content = result.document.export_to_markdown() if self.preserve_tables else result.document.export_to_text()

            metadata = {
                "title": file_path.stem,
                "file_type": "pdf",
                "page_count": len(result.document.pages) if hasattr(result.document, "pages") else 0,
                "has_tables": bool(result.document.tables) if hasattr(result.document, "tables") else False,
                "structure": self._extract_structure(result),
            }

            if self.preserve_tables and hasattr(result.document, "tables"):
                metadata["tables"] = [self._extract_table(table) for table in result.document.tables]

            return {"content": content, "metadata": metadata}

        except ImportError:
            logger.warning("Docling not available, falling back to basic PDF extraction")
            return await self._process_pdf_fallback(file_path)
        except Exception as e:
            logger.error(f"Docling PDF processing failed: {e}", exc_info=True)
            return await self._process_pdf_fallback(file_path)

    async def _process_pdf_fallback(self, file_path: Path) -> Dict[str, Any]:
        """Fallback PDF processor using PyPDF2."""
        try:
            import PyPDF2

            content = []
            with open(file_path, "rb") as file:
                pdf_reader = PyPDF2.PdfReader(file)
                page_count = len(pdf_reader.pages)

                for page in pdf_reader.pages:
                    content.append(page.extract_text())

            return {
                "content": "\n\n".join(content),
                "metadata": {
                    "title": file_path.stem,
                    "file_type": "pdf",
                    "page_count": page_count,
                    "has_tables": False,
                    "processing_method": "fallback",
                },
            }
        except Exception as e:
            logger.error(f"Fallback PDF processing failed: {e}")
            raise ValueError(f"Could not process PDF: {e}")

    async def _process_docx(self, file_path: Path) -> Dict[str, Any]:
        """Process DOCX using Docling."""
        try:
            from docling.document_converter import DocumentConverter

            converter = DocumentConverter()
            result = converter.convert(str(file_path))

            content = result.document.export_to_markdown() if self.preserve_tables else result.document.export_to_text()

            metadata = {
                "title": file_path.stem,
                "file_type": "docx",
                "has_tables": bool(result.document.tables) if hasattr(result.document, "tables") else False,
                "structure": self._extract_structure(result),
            }

            return {"content": content, "metadata": metadata}

        except ImportError:
            logger.warning("Docling not available, falling back to basic DOCX extraction")
            return await self._process_docx_fallback(file_path)
        except Exception as e:
            logger.error(f"Docling DOCX processing failed: {e}", exc_info=True)
            return await self._process_docx_fallback(file_path)

    async def _process_docx_fallback(self, file_path: Path) -> Dict[str, Any]:
        """Fallback DOCX processor using python-docx."""
        try:
            from docx import Document

            doc = Document(file_path)
            content = "\n\n".join([para.text for para in doc.paragraphs if para.text.strip()])

            return {
                "content": content,
                "metadata": {
                    "title": file_path.stem,
                    "file_type": "docx",
                    "has_tables": len(doc.tables) > 0,
                    "processing_method": "fallback",
                },
            }
        except Exception as e:
            logger.error(f"Fallback DOCX processing failed: {e}")
            raise ValueError(f"Could not process DOCX: {e}")

    async def _process_excel(self, file_path: Path) -> Dict[str, Any]:
        """Process Excel files."""
        try:
            import pandas as pd

            # Read all sheets
            excel_file = pd.ExcelFile(file_path)
            content_parts = []

            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                content_parts.append(f"## Sheet: {sheet_name}\n\n{df.to_markdown(index=False)}")

            return {
                "content": "\n\n".join(content_parts),
                "metadata": {
                    "title": file_path.stem,
                    "file_type": "xlsx",
                    "sheet_count": len(excel_file.sheet_names),
                    "has_tables": True,
                },
            }
        except Exception as e:
            logger.error(f"Excel processing failed: {e}")
            raise ValueError(f"Could not process Excel file: {e}")

    async def _process_ppt(self, file_path: Path) -> Dict[str, Any]:
        """Process PowerPoint files."""
        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            content_parts = []

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)

                if slide_text:
                    content_parts.append(f"## Slide {slide_num}\n\n" + "\n".join(slide_text))

            return {
                "content": "\n\n".join(content_parts),
                "metadata": {
                    "title": file_path.stem,
                    "file_type": "pptx",
                    "slide_count": len(prs.slides),
                },
            }
        except Exception as e:
            logger.error(f"PowerPoint processing failed: {e}")
            raise ValueError(f"Could not process PowerPoint file: {e}")

    async def _process_text(self, file_path: Path) -> Dict[str, Any]:
        """Process plain text files."""
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                content = f.read()

            return {
                "content": content,
                "metadata": {
                    "title": file_path.stem,
                    "file_type": file_path.suffix.lstrip("."),
                },
            }
        except Exception as e:
            logger.error(f"Text file processing failed: {e}")
            raise ValueError(f"Could not process text file: {e}")

    def _extract_structure(self, result: Any) -> Dict[str, Any]:
        """Extract document structure from Docling result."""
        try:
            return {
                "has_headings": bool(getattr(result.document, "headings", [])),
                "has_lists": bool(getattr(result.document, "lists", [])),
                "has_images": bool(getattr(result.document, "images", [])),
            }
        except:
            return {}

    def _extract_table(self, table: Any) -> Dict[str, Any]:
        """Extract table data from Docling table object."""
        try:
            return {
                "rows": len(table.rows) if hasattr(table, "rows") else 0,
                "columns": len(table.columns) if hasattr(table, "columns") else 0,
                "data": table.to_dict() if hasattr(table, "to_dict") else {},
            }
        except:
            return {}


async def process_document(file_path: Path) -> Dict[str, Any]:
    """
    Convenience function to process a document file.

    Args:
        file_path: Path to the document file

    Returns:
        Dictionary with content and metadata
    """
    processor = DoclingProcessor()
    return await processor.process_file(file_path)
